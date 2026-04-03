"""Ingest raw content into the KB raw/ directory.

Supports: .md/.txt (copy), .docx (mammoth), .pptx (python-pptx),
          .pdf (pdfminer), URLs (html2text), images (copy).
Each ingested file gets a .meta.json sidecar with source info + domain override.
"""

import contextlib
import json
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from urllib.parse import urlparse
from urllib.request import urlopen

from brainery.config import get_kb_path


def run(args, cfg):
    """Ingest raw content from a file or URL into the KB."""
    kb = args.kb or cfg.get("default_kb", "personal")
    kb_path = get_kb_path(cfg, kb)
    raw_dir = kb_path / "raw"
    raw_dir.mkdir(parents=True, exist_ok=True)

    source = args.source
    domain_override = getattr(args, "domain", None)

    if source.startswith("http://") or source.startswith("https://"):
        filename, content = _ingest_url(source)
    else:
        source_path = Path(source).expanduser()
        if not source_path.exists():
            print(f"[error] File not found: {source}")
            sys.exit(1)
        filename, content = _ingest_file(source_path)

    # Save the raw content
    dest_file = raw_dir / filename
    dest_file.write_text(content, encoding="utf-8")

    # Save metadata sidecar
    meta = {
        "source": source,
        "domain_override": domain_override,
        "ingested": datetime.utcnow().isoformat() + "Z",
        "source_type": "url" if source.startswith("http") else "file",
    }
    meta_file = raw_dir / f"{filename}.meta.json"
    meta_file.write_text(json.dumps(meta, indent=2), encoding="utf-8")

    print(f"✓ Ingested: {filename}")
    print("Run 'kb compile' to add to your wiki.")


def _ingest_file(source_path: Path) -> tuple[str, str]:
    """Ingest a local file, converting as needed. Returns (filename, content)."""
    suffix = source_path.suffix.lower()
    timestamp = datetime.utcnow().isoformat() + "Z"
    header = f"<!-- Source: {source_path.name} -->\n<!-- Ingested: {timestamp} -->\n\n"

    if suffix == ".md" or suffix == ".txt":
        content = source_path.read_text(encoding="utf-8")
        return (source_path.name, header + content)

    elif suffix == ".docx":
        content = _convert_docx(source_path)
        filename = source_path.stem + ".md"
        return (filename, header + content)

    elif suffix == ".pptx":
        content = _convert_pptx(source_path)
        filename = source_path.stem + ".md"
        return (filename, header + content)

    elif suffix == ".pdf":
        content = _convert_pdf(source_path)
        filename = source_path.stem + ".md"
        return (filename, header + content)

    else:
        # Copy as-is for images and other binary types
        content = source_path.read_bytes()
        return (source_path.name, content)


def _ingest_url(url: str) -> tuple[str, str]:
    """Ingest a URL using html2text. Returns (filename, content)."""
    try:
        _ensure_package("html2text")
        import html2text
    except ImportError:
        print("[error] html2text not found. Install: pip install html2text")
        sys.exit(1)

    try:
        with urlopen(url) as response:
            html_content = response.read().decode("utf-8")
    except Exception as e:
        print(f"[error] Failed to fetch URL: {e}")
        sys.exit(1)

    h = html2text.HTML2Text()
    h.ignore_links = False
    h.body_width = 0
    markdown_content = h.handle(html_content)

    # Derive filename from URL slug
    parsed = urlparse(url)
    slug = parsed.path.strip("/").split("/")[-1] or parsed.netloc
    slug = slug.replace(".", "-").replace("?", "")[:80]
    if not slug.endswith(".md"):
        slug += ".md"

    timestamp = datetime.utcnow().isoformat() + "Z"
    header = f"<!-- Source URL: {url} -->\n<!-- Ingested: {timestamp} -->\n\n"
    return (slug, header + markdown_content)


def _convert_docx(source_path: Path) -> str:
    """Convert .docx to Markdown using mammoth."""
    try:
        _ensure_package("mammoth")
        import mammoth
    except ImportError:
        print("[error] mammoth not found. Install: pip install mammoth")
        sys.exit(1)

    try:
        with open(source_path, "rb") as f:
            result = mammoth.convert_docx(f)
            return result.value
    except Exception as e:
        print(f"[error] Failed to convert docx: {e}")
        sys.exit(1)


def _convert_pptx(source_path: Path) -> str:
    """Convert .pptx to Markdown by extracting slide text."""
    try:
        _ensure_package("python-pptx")
        from pptx import Presentation
    except ImportError:
        print("[error] python-pptx not found. Install: pip install python-pptx")
        sys.exit(1)

    try:
        prs = Presentation(source_path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"## Slide {i}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
            slides_text.append(slide_text)
        return "\n".join(slides_text)
    except Exception as e:
        print(f"[error] Failed to convert pptx: {e}")
        sys.exit(1)


def _convert_pdf(source_path: Path) -> str:
    """Convert .pdf to text using pdfminer."""
    try:
        _ensure_package("pdfminer.six")
        from pdfminer.high_level import extract_text
    except ImportError:
        print("[error] pdfminer.six not found. Install: pip install pdfminer.six")
        sys.exit(1)

    try:
        text = extract_text(str(source_path))
        return text
    except Exception as e:
        print(f"[error] Failed to convert pdf: {e}")
        sys.exit(1)


def _ensure_package(package_name: str) -> None:
    """Auto-install package via pip if missing."""
    try:
        __import__(package_name)
    except ImportError:
        # Lazy install
        with contextlib.suppress(subprocess.CalledProcessError):
            subprocess.run(
                [sys.executable, "-m", "pip", "install", "--break-system-packages", "-q", package_name],
                check=True,
            )
